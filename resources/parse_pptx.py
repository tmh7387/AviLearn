import sys
import json
import os

try:
    from pptx import Presentation
except ImportError:
    print(json.dumps({
        "error": "Missing dependency 'python-pptx'. Please install it using 'pip install python-pptx'."
    }))
    sys.exit(1)

def parse_pptx(file_path):
    if not os.path.exists(file_path):
        print(json.dumps({"error": f"File not found: {file_path}"}))
        sys.exit(1)

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(json.dumps({"error": f"Failed to open presentation: {str(e)}"}))
        sys.exit(1)

    slides_data = []

    for index, slide in enumerate(prs.slides):
        title = ""
        text_content = []
        has_images = False
        has_charts = False

        # Try to find a title shape or guess title
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            # Check for text frame
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != title:
                    text_content.append(text)
            
            # Check for image type
            # shape_type 13 is picture
            if hasattr(shape, 'shape_type'):
                if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                    has_images = True
                elif shape.shape_type == 3: # MSO_SHAPE_TYPE.CHART
                    has_charts = True

        # Extract notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append({
            "index": index + 1,
            "title": title or f"Slide {index + 1}",
            "textContent": "\n".join(text_content),
            "hasImages": has_images,
            "hasCharts": has_charts,
            "notes": notes
        })

    print(json.dumps({"slides": slides_data}))

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: python parse_pptx.py <path_to_pptx>"}))
        sys.exit(1)
    
    parse_pptx(sys.argv[1])
